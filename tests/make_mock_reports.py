#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""make_mock_reports.py — 生成植入已知缺陷的模拟 DV 报告（.docx / .pptx），用于验证脚本。

开发期工具，需要 python-docx 与 python-pptx（skill 本身不需要）。
缺陷植入清单见文件末尾 EXPECTED，供验证时逐条核对。

用法：
  python make_mock_reports.py [输出目录]   # 默认输出到本脚本所在目录
"""
import base64
import os
import sys

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PPt

# 1x1 像素 PNG，用于 PPT 图片页
PNG_1PX = base64.b64decode(
    'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==')

SAMPLE_ID = 'BTS-2024-0187'
REPORT_ID = 'BTR-DV-2024-0156'

SUMMARY_ROWS = [
    ('1', '过放电保护', '6.1.1', '合格'),
    ('2', '过充电保护', '6.1.2', '合格'),
    ('3', '短路保护', '6.1.3', '合格'),
    ('4', '振动', '6.2.1', '不合格'),          # ← 与结论"全部合格"矛盾（CD-S01）
    ('5', '机械冲击', '6.2.2', '合格'),
    ('6', '湿热循环', '6.3.1', ''),             # ← 空白判定单元格（CD-G07）
    ('7', '温度冲击', '6.3.2', '合格'),
    ('8', '外部火烧', '6.4.1', '待补充'),        # ← 占位符（CD-G07）
]


def build_docx(out_dir):
    doc = Document()
    doc.add_heading('动力电池包 DV 测试报告', 0)
    for line in [f'报告编号：{REPORT_ID}', '委托单位：某新能源科技有限公司',
                 '生产单位：某电池系统有限公司', '样品名称：动力电池包',
                 '型号规格：PACK-60，60 kwh',          # ← 单位写法 kwh（CD-J01）
                 '报告日期：2024年3月15日']:
        doc.add_paragraph(line)

    doc.add_heading('签署', 1)
    t = doc.add_table(4, 2)
    for r, (k, v) in enumerate([('编制', '张三 2024年3月14日'), ('校对', '李四 2024年3月14日'),
                                ('审核', '王五 2024年3月15日'), ('批准', '赵六 2024年3月15日')]):
        t.cell(r, 0).text, t.cell(r, 1).text = k, v

    doc.add_heading('1 概述', 1)
    doc.add_paragraph('委托日期：2024年3月1日。')
    doc.add_paragraph('试验时间：2024年2月20日 至 2024年3月10日。')  # ← 倒挂（CD-S08）
    doc.add_paragraph('本报告依据客户 DV 大纲及 GB 38031-2015 编制。')  # ← 作废年号（CD-G01）

    doc.add_heading('2 样品信息', 1)
    doc.add_paragraph(f'样品编号：{SAMPLE_ID}，数量：2 台，全新样品。')
    doc.add_paragraph('BMS 硬件版本：V2.1；软件版本：V3.4.5。生产批次：2024-01。')
    doc.add_paragraph(f'其中 1 台样品（{SAMPLE_ID.replace("0187", "01B7")}）用于振动试验。')  # ← 疑似笔误（CD-S04）

    doc.add_heading('3 引用标准', 1)
    doc.add_paragraph('GB 38031-2020《电动汽车用动力蓄电池安全要求》')  # ← 年号前后不一致（CD-G02）
    doc.add_paragraph('GB/T 31486-2015《电动汽车用动力蓄电池电性能要求及试验方法》')

    doc.add_heading('4 试验设备', 1)
    t = doc.add_table(2, 4)
    for c, h in enumerate(['设备名称', '型号', '编号', '校准有效期至']):
        t.cell(0, c).text = h
    for c, v in enumerate(['充放电测试系统', 'CT-4004', 'EQ-0098', '2024年1月31日']):  # ← 过期（CD-S03）
        t.cell(1, c).text = v

    doc.add_heading('5 试验项目汇总', 1)
    t = doc.add_table(len(SUMMARY_ROWS) + 1, 4)
    for c, h in enumerate(['序号', '试验项目', '依据条款', '判定']):
        t.cell(0, c).text = h
    for r, row in enumerate(SUMMARY_ROWS, 1):
        for c, v in enumerate(row):
            t.cell(r, c).text = v

    doc.add_heading('6 单项试验结果', 1)
    doc.add_heading('6.1 振动试验', 2)
    doc.add_paragraph(f'样品编号：{SAMPLE_ID.replace("-", "_")}。')  # ← 写法不一致
    doc.add_paragraph('试验后绝缘阻值 38.2MΩ（限值 ≥100Ω/V），判定：合格。')
    doc.add_paragraph('试验后采样内阻 38.2mΩ（限值 ≤35mΩ），判定：合格。')  # ← 超限判合格（DL-A05）
    doc.add_heading('6.2 过放电保护试验', 2)
    doc.add_paragraph('环境温度 25 °C，实测保护动作电压 33.9V（限值 ≤35V 判定合格）。')  # ← 临界值 + °C
    doc.add_heading('6.3 湿热循环试验', 2)
    doc.add_paragraph('试验温度上限 55 ℃，湿度 95%RH。')  # ← ℃ 与 °C 混用（CD-J01）
    doc.add_paragraph('试验后样品外观：待补充。')  # ← 占位符

    doc.add_heading('7 结论', 1)
    doc.add_paragraph('本次共 8 项试验，8 项全部合格。结论：合格。')  # ← 与汇总表矛盾（CD-S01）
    doc.add_paragraph('整改措施：××。')  # ← 占位符（CD-G07）

    path = os.path.join(out_dir, 'mock_dv_report.docx')
    doc.save(path)
    return path


def _textbox(slide, left, top, width, height, lines, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = PPt(size)
    return box


def build_pptx(out_dir):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # S1 封面
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.5), Inches(9), Inches(2),
             ['动力电池包 DV 测试汇报', f'报告编号：{REPORT_ID}', '2024年3月'], 24)

    # S2 汇总表（含 1 不合格 + 1 空格）
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8), ['试验项目汇总'], 24)
    rows, cols = len(SUMMARY_ROWS) + 1, 4
    gt = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2),
                            Inches(9), Inches(5)).table
    for c, h in enumerate(['序号', '试验项目', '依据条款', '判定']):
        gt.cell(0, c).text = h
    for r, row in enumerate(SUMMARY_ROWS, 1):
        for c, v in enumerate(row):
            gt.cell(r, c).text = v

    # S3 关键结果（文字与表矛盾 + 无图表来源 + kwh）
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.5), Inches(9), Inches(4),
             ['关键结果', '全部 8 项试验均合格', '容量保持率 96.5%',
              '电池包能量 60 kwh', '湿热试验上限 55 ℃'])

    # S4 纯图片页
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8), ['试验现场照片'], 24)
    png = os.path.join(out_dir, '_px.png')
    with open(png, 'wb') as f:
        f.write(PNG_1PX)
    s.shapes.add_picture(png, Inches(1), Inches(1.5), Inches(3), Inches(3))

    # S5 问题与整改（占位符）
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.5), Inches(9), Inches(3),
             ['问题与整改', '振动试验整改措施：待补充'])

    # S6 结论（含备注）
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.5), Inches(0.5), Inches(9), Inches(3),
             ['结论', '结论：合格', '8 项全部合格'])
    s.notes_slide.notes_text_frame.text = '汇报人：张三；汇报日期 2024年3月20日'

    path = os.path.join(out_dir, 'mock_dv_report.pptx')
    prs.save(path)
    os.remove(png)
    return path


EXPECTED = """
脚本应命中的植入缺陷：
[docx] CHECK-1 待补充/××；CHECK-2 委托晚于试验开始、校准有效期早于试验结束；
       CHECK-3 BTS-2024-01B7 疑似笔误、BTS_2024_0187 写法不一致；
       CHECK-4 kwh、℃/°C 混用；CHECK-5 不合格1项但文字称8项全部合格且结论合格；
       CHECK-6 GB 38031 两个年号不一致+2015 疑似作废、GB/T 31486 提示；CHECK-7 空白判定格。
[pptx] 提取标记第4页图片需人工核对；CHECK-1 待补充；CHECK-4 kwh；
       CHECK-5 不合格1项但文字称全部合格；CHECK-7 空白判定格。
语义审核应额外命中：6.1 内阻超限判合格（DL-A05）、33.9V 临界值（DL-A03）、
       PPT 图表无来源（DC-P03）。
"""

if __name__ == '__main__':
    out_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    os.makedirs(out_dir, exist_ok=True)
    p1 = build_docx(out_dir)
    p2 = build_pptx(out_dir)
    print(f'OK: {p1}')
    print(f'OK: {p2}')
    print(EXPECTED)
