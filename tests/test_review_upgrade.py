#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Regression tests for the evidence-first review upgrade."""

import importlib.util
import base64
import io
import json
import os
import tempfile
import unittest


ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def load_script(name):
    path = os.path.join(ROOT, "scripts", name)
    spec = importlib.util.spec_from_file_location(name.replace(".", "_"), path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


LANGUAGE_REVIEW = load_script("language_review.py")
FIGURE_CHECKS = load_script("figure_checks.py")
NORMALIZE_REVIEW = load_script("normalize_review.py")
EVIDENCE_PIPELINE = load_script("evidence_pipeline.py")
RENDER_EVIDENCE = load_script("render_evidence.py")
RENDER_REVIEW = load_script("render_review.py")
LOCATION_RESOLVER = load_script("location_resolver.py")


def sample_evidence():
    return {
        "schema_version": "evidence.v1",
        "document": {
            "source_path": r"C:\reports\demo.docx",
            "source_name": "demo.docx",
            "source_type": "docx",
        },
        "units": [
            {
                "order": 1,
                "anchor": "[P0001]",
                "kind": "paragraph",
                "text": "4.2 Capacity test: display capacity variation from 0 to 20000 s.",
                "section_context": ["4.2 Capacity test"],
                "neighbor_previous": None,
                "neighbor_next": {
                    "anchor": "[P0002]",
                    "text": "The result is shown in Figure 4-3.",
                },
                "direct_figure_references": [],
                "caption_candidates": [],
            },
            {
                "order": 2,
                "anchor": "[P0002]",
                "kind": "paragraph",
                "text": "The result is shown in Figure 4-3.",
                "section_context": ["4.2 Capacity test"],
                "neighbor_previous": {
                    "anchor": "[P0001]",
                    "text": "4.2 Capacity test: display capacity variation from 0 to 20000 s.",
                },
                "neighbor_next": None,
                "direct_figure_references": [
                    {"normalized_id": "4-3", "raw": "Figure 4-3"}
                ],
                "caption_candidates": [],
            },
        ],
        "actual_figures": [
            {
                "id": "AF0001",
                "kind": "image",
                "figure_id": "Figure 4-3",
                "source_anchor": "[P0003]",
                "extracted_path": r"media\AF0001_image.png",
                "sha256": "same-image",
                "excluded": False,
                "matched_expected_ids": ["EF0001"],
                "caption_candidates": ["Figure 4-3"],
                "nearby_text": {},
            },
            {
                "id": "AF0002",
                "kind": "image",
                "figure_id": "Figure 4-4",
                "source_anchor": "[P0004]",
                "extracted_path": r"media\AF0002_image.png",
                "sha256": "same-image",
                "excluded": False,
                "matched_expected_ids": [],
                "caption_candidates": ["Figure 4-4"],
                "nearby_text": {},
            },
        ],
        "expected_figures": [
            {
                "expected_id": "EF0001",
                "figure_id": "Figure 4-3",
                "source_types": ["direct_reference"],
                "source_anchors": ["[P0002]"],
                "excerpts": ["capacity variation from 0 to 20000 s"],
                "evidence": [{"anchor": "[P0001]", "excerpt": "0 to 20000 s"}],
                "keywords": ["capacity", "Ah"],
                "axis_unit_range_hints": {
                    "axis": ["capacity", "time"],
                    "units": ["Ah", "s"],
                    "ranges": ["0-20000 s"],
                },
                "required": True,
            }
        ],
        "matches": [
            {
                "expected_id": "EF0001",
                "actual_ids": ["AF0001"],
                "status": "matched",
                "score": 100,
            }
        ],
        "extraction_warnings": [],
    }


class ReviewUpgradeTests(unittest.TestCase):
    def test_docx_media_and_figure_context_extraction(self):
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx is not installed")

        png = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
            "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
        )
        with tempfile.TemporaryDirectory() as temp_dir:
            source = os.path.join(temp_dir, "figure-demo.docx")
            artifact_dir = os.path.join(temp_dir, "figure-demo.review-artifacts")
            doc = Document()
            doc.add_heading("4.2 Capacity Test", level=1)
            doc.add_paragraph(
                "The capacity variation from 0 to 20000 s is shown in Figure 4-3."
            )
            doc.add_picture(io.BytesIO(png))
            doc.add_paragraph("Figure 4-3 Capacity variation during cycling")
            doc.save(source)

            result = EVIDENCE_PIPELINE.build_evidence(source, artifact_dir)
            self.assertEqual(result["document"]["actual_figure_count"], 1)
            self.assertTrue(result["actual_figures"][0]["extracted_path"])
            self.assertTrue(
                any(item.get("normalized_figure_id") == "4-3"
                    for item in result["expected_figures"])
            )
            self.assertEqual(result["units"][0]["location_detail"]["source_anchor"], "[P0001]")
            self.assertEqual(result["units"][0]["location_detail"]["page_status"], "unavailable")
            self.assertIn("页码未解析", result["units"][0]["location_detail"]["display"])
            self.assertIn("location-P0001", RENDER_EVIDENCE.render_evidence_html(result))
            extracted = os.path.join(
                artifact_dir,
                result["actual_figures"][0]["extracted_path"],
            )
            self.assertTrue(os.path.exists(extracted))

    def test_language_input_preserves_context_and_terms(self):
        result = LANGUAGE_REVIEW.build_review_input(sample_evidence())
        self.assertEqual(result["schema_version"], "language-review-input.v1")
        self.assertEqual(len(result["units"]), 2)
        self.assertEqual(result["units"][0]["section_context"], ["4.2 Capacity test"])
        self.assertIn("Capacity", result["units"][0]["technical_terms"])
        self.assertEqual(result["units"][1]["nearby_figures"], [])

    def test_figure_checks_no_vision_is_explicit(self):
        result = FIGURE_CHECKS.build_figure_review(
            sample_evidence(), observations=[], vision_available=False
        )
        categories = {item["category"] for item in result["issues"]}
        self.assertIn("figure_vision_unavailable", categories)
        self.assertIn("figure_exact_duplicate", categories)
        self.assertTrue(
            all(item["requires_author_confirmation"] for item in result["issues"])
        )

    def test_vision_observations_detect_variable_and_range_mismatch(self):
        result = FIGURE_CHECKS.build_figure_review(
            sample_evidence(),
            observations=[
                {
                    "actual_id": "AF0001",
                    "confidence": "high",
                    "x_axis": "Time/s",
                    "y_axis": "Voltage/V",
                    "units": ["V", "s"],
                    "x_range": "0-200 s",
                    "trend": "increasing",
                }
            ],
            vision_available=True,
        )
        categories = {item["category"] for item in result["issues"]}
        self.assertIn("figure_context_mismatch", categories)
        self.assertIn("figure_data_range_mismatch", categories)

    def test_normalize_parses_numeric_confidence(self):
        raw = {
            "metadata": {"report_name": "demo"},
            "issues": [
                {
                    "issue_id": "ISS-001",
                    "category": "figure_context_mismatch",
                    "severity": "critical",
                    "confidence": "0.97",
                    "location": "[P0001]",
                    "expected": "Capacity/Ah",
                    "actual": "Voltage/V",
                    "evidence": "正文要求容量图，图中为电压图",
                    "recommendation": "替换图片",
                }
            ],
        }
        result = NORMALIZE_REVIEW.normalize(raw)
        self.assertEqual(result["issues"][0]["confidence"], "高")
        self.assertEqual(result["validation_warnings"], [])
        self.assertEqual(result["summary"]["severity_counts"]["严重"], 1)

    def test_location_page_map_resolver_uses_explicit_pages(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            page_map = os.path.join(temp_dir, "pages.json")
            with open(page_map, "w", encoding="utf-8") as handle:
                json.dump({"pages": {"[P0001]": 4, "[T01]": 5}}, handle)
            result = LOCATION_RESOLVER.resolve_docx_pages(
                os.path.join(temp_dir, "demo.docx"),
                [{"anchor": "[P0001]"}, {"anchor": "[T01]"}],
                page_map_path=page_map,
            )
            self.assertEqual(result["status"], "resolved")
            self.assertEqual(result["pages"]["[P0001]"], 4)
            self.assertEqual(result["pages"]["[T01]"], 5)

    def test_pptx_location_includes_slide_shape_and_geometry(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx is not installed")

        with tempfile.TemporaryDirectory() as temp_dir:
            source = os.path.join(temp_dir, "location-demo.pptx")
            artifact_dir = os.path.join(temp_dir, "location-demo.review-artifacts")
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(3), Inches(1)
            ).text_frame.text = "图表说明"
            prs.save(source)

            result = EVIDENCE_PIPELINE.build_evidence(source, artifact_dir)
            unit = next(item for item in result["units"] if item["anchor"] == "[S01-1]")
            detail = unit["location_detail"]
            self.assertEqual(detail["slide"], 1)
            self.assertEqual(detail["shape"], 1)
            self.assertEqual(detail["object_type"], "shape")
            self.assertEqual(detail["position"]["unit"], "EMU")
            self.assertIn("左侧", detail["position_summary"])
            self.assertIn("第 1 页", detail["display"])
            self.assertIn("第 1 个形状", detail["display"])

    def test_normalize_enriches_legacy_location_from_evidence(self):
        evidence = sample_evidence()
        evidence["units"][0]["location_detail"] = {
            "schema_version": "location.v1",
            "display": "页码未解析 · 第 1 个段落 · 章节：4.2 Capacity test · [P0001]",
            "source_anchor": "[P0001]",
            "page_status": "unavailable",
            "evidence_path": r"C:\reports\demo.review-artifacts\evidence.html",
            "evidence_anchor": "location-P0001",
        }
        result = NORMALIZE_REVIEW.normalize(
            {"issues": [{
                "issue_id": "ISS-LOC",
                "category": "language",
                "severity": "major",
                "confidence": "high",
                "location": "[P0001]",
                "expected": "清晰表达",
                "actual": "存在歧义",
                "evidence": "原文摘录",
                "recommendation": "改写",
            }]},
            evidence=evidence,
        )
        issue = result["issues"][0]
        self.assertEqual(issue["location_display"], evidence["units"][0]["location_detail"]["display"])
        self.assertEqual(issue["location_detail"]["evidence_anchor"], "location-P0001")

    def test_render_location_has_source_and_evidence_links(self):
        data = {
        "metadata": {"report_name": "location demo"},
        "summary": {},
        "issues": [{
            "issue_id": "ISS-LOC",
            "category": "定位",
            "severity": "一般",
            "confidence": "高",
            "location": "[S08-3]",
            "location_detail": {
                "display": "第 8 页 · 第 3 个形状（图表） · [S08-3]",
                "source_anchor": "[S08-3]",
                "source_file": r"C:\reports\demo.pptx",
                "evidence_path": r"C:\reports\demo.review-artifacts\evidence.html",
                "evidence_anchor": "location-S08-3",
            },
            "expected": "图表位置清晰",
            "actual": "需要人工复核",
            "evidence": "位置证据",
            "recommendation": "打开证据定位",
        }],
        }
        html_text = RENDER_REVIEW.render_html(data, "location demo")
        self.assertIn("第 8 页", html_text)
        self.assertIn("打开原文件", html_text)
        self.assertIn("打开证据定位", html_text)
        self.assertIn("file:///C:/reports/demo.pptx", html_text)
        self.assertIn("#location-S08-3", html_text)


if __name__ == "__main__":
    unittest.main()
